# -*- coding: utf-8 -*-
"""
Apresentação IMPACTO — dados e alertas reais do SIS Clima-Saúde MT.

Lê Postgres/SQLite operacional, gera alertas multinível ao vivo e monta PPTX.
Uso:
  .venv\\Scripts\\python.exe gerar_apresentacao_sis_impacto_real.py
  .venv\\Scripts\\python.exe gerar_apresentacao_sis_impacto_real.py --also-downloads
"""
from __future__ import annotations

import argparse
import json
import shutil
from datetime import datetime
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "docs" / "apresentacoes"
DEFAULT_OUT = OUT_DIR / "Apresentacao_SIS_Clima_Saude_IMPACTO_dados_reais.pptx"

GREEN_DARK = RGBColor(0x08, 0x35, 0x2E)
GREEN = RGBColor(0x0F, 0x6E, 0x56)
GREEN_SOFT = RGBColor(0x1A, 0x8A, 0x6E)
CREAM = RGBColor(0xF3, 0xF6, 0xF4)
INK = RGBColor(0x1A, 0x2E, 0x28)
MUTED = RGBColor(0x4B, 0x5C, 0x56)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
PURPLE = RGBColor(0x6D, 0x28, 0xD9)
YELLOW = RGBColor(0xC4, 0x92, 0x00)

LEVEL_RGB = {
    "verde": GREEN_SOFT,
    "amarela": YELLOW,
    "laranja": ORANGE,
    "vermelha": RED,
    "roxa": PURPLE,
    "cinza": MUTED,
}


def _set_run(run, *, size=18, bold=False, color=INK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _add_bg(slide, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def _textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color)
    return box


def _bullets(slide, left, top, width, height, items, *, size=15, color=INK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"• {item}"
        _set_run(run, size=size, color=color)
    return box


def _header_bar(slide, title, subtitle=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN_DARK
    bar.line.fill.background()
    _textbox(slide, 0.5, 0.18, 12, 0.4, title, size=24, bold=True, color=WHITE)
    if subtitle:
        _textbox(slide, 0.5, 0.58, 12, 0.35, subtitle, size=13, color=RGBColor(0xD7, 0xEB, 0xE3))


def _footer(slide, n, total):
    _textbox(slide, 0.5, 7.1, 10, 0.3, "SIS Clima-Saúde MT · CIEVS/SES-MT · dados reais do painel", size=10, color=MUTED)
    _textbox(slide, 11.2, 7.1, 1.6, 0.3, f"{n}/{total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def _card(slide, left, top, width, height, title, body, accent=GREEN):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = RGBColor(0xD5, 0xE0, 0xDA)
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.12), Inches(height))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()
    _textbox(slide, left + 0.25, top + 0.12, width - 0.35, 0.32, title, size=13, bold=True, color=GREEN_DARK)
    _textbox(slide, left + 0.25, top + 0.45, width - 0.35, height - 0.55, body, size=12, color=INK)


def _fmt_num(v, nd=1, suffix=""):
    try:
        if v is None or (isinstance(v, float) and pd.isna(v)):
            return "—"
        x = float(v)
        if abs(x) >= 1000:
            return f"{x:,.0f}{suffix}".replace(",", ".")
        return f"{x:.{nd}f}{suffix}"
    except Exception:
        return str(v) if v is not None else "—"


def load_real_bundle() -> dict:
    from sisclima.core.db import backend_name, read_table
    from sisclima.engines.alertas_multinivel import build_alertas_multinivel
    from sisclima.engines.indice_pressao_saude import (
        build_indice_pressao_municipal,
        state_pressao_summary,
    )

    resumo = read_table("resumo_municipal_atual")
    pred = read_table("predicao_calor_7d_municipal_v6")
    sis = read_table("ops_sisreg_municipio")
    alerta_int = read_table("alerta_integrado_sis_titan")
    sim = read_table("sim_obitos_calor_municipal_v6")
    saude = read_table("saude_calor_municipio")

    if resumo is None or resumo.empty:
        raise RuntimeError("resumo_municipal_atual vazio — suba o Docker/DB antes.")

    press = build_indice_pressao_municipal(
        resumo,
        sim_mun=sim if not sim.empty else None,
        saude_calor_mun=saude if not saude.empty else None,
        pred_7d=pred if not pred.empty else None,
        sisreg=sis if not sis.empty else None,
    )
    payloads = build_alertas_multinivel(
        resumo,
        alerta_integrado=alerta_int if not alerta_int.empty else None,
        predicao_7d=pred if not pred.empty else None,
        min_level="laranja",
    )

    niv = resumo["nivel"].astype(str).str.lower() if "nivel" in resumo.columns else pd.Series(dtype=str)
    dist = {k: int((niv == k).sum()) for k in ("verde", "amarela", "laranja", "vermelha", "roxa", "cinza")}
    top = resumo.sort_values(
        [c for c in ("score", "risco_cumulativo_3d", "tmax") if c in resumo.columns],
        ascending=False,
    ).head(12)

    est = next((p for p in payloads if p.get("escopo") == "estadual"), None)
    cui = next((p for p in payloads if p.get("escopo") == "cuiaba"), None)
    regs = sorted(
        [p for p in payloads if p.get("escopo") == "regional"],
        key=lambda p: {"roxa": 4, "vermelha": 3, "laranja": 2, "amarela": 1}.get(str(p.get("nivel")), 0),
        reverse=True,
    )
    muns = sorted(
        [p for p in payloads if p.get("escopo") == "municipal"],
        key=lambda p: {"roxa": 4, "vermelha": 3, "laranja": 2, "amarela": 1}.get(str(p.get("nivel")), 0),
        reverse=True,
    )

    sis_top = pd.DataFrame()
    if not sis.empty and "solicitacoes_abertas" in sis.columns:
        sis_top = sis.sort_values("solicitacoes_abertas", ascending=False).head(8)

    return {
        "gerado_em": datetime.now().isoformat(timespec="seconds"),
        "backend": backend_name(),
        "n_mun": int(resumo["cod_ibge"].nunique()) if "cod_ibge" in resumo.columns else len(resumo),
        "dist": dist,
        "tmax_max": float(pd.to_numeric(resumo.get("tmax"), errors="coerce").max()) if "tmax" in resumo.columns else None,
        "tmax_med": float(pd.to_numeric(resumo.get("tmax"), errors="coerce").mean()) if "tmax" in resumo.columns else None,
        "utci_max": float(pd.to_numeric(resumo.get("utci_proxy"), errors="coerce").max()) if "utci_proxy" in resumo.columns else None,
        "ocup_med": float(pd.to_numeric(resumo.get("ocupacao_leitos_pct"), errors="coerce").mean()) if "ocupacao_leitos_pct" in resumo.columns else None,
        "pressao_state": state_pressao_summary(press),
        "sisreg_n": int(len(sis)),
        "sisreg_fonte": str(sis["fonte"].dropna().iloc[0]) if not sis.empty and "fonte" in sis.columns and sis["fonte"].notna().any() else "—",
        "sisreg_sols_total": float(pd.to_numeric(sis.get("solicitacoes_abertas"), errors="coerce").sum()) if not sis.empty else None,
        "sisreg_top": sis_top,
        "top_mun": top,
        "payloads": payloads,
        "est": est,
        "cui": cui,
        "regs": regs,
        "muns": muns,
        "n_alertas": len(payloads),
    }


def _inds_lines(inds, n=6) -> list[str]:
    out = []
    for item in (inds or [])[:n]:
        out.append(f"{item.get('rotulo')}: {item.get('valor')}")
    return out or ["Sem indicadores nesta rodada."]


def build(out: Path) -> tuple[Path, dict]:
    data = load_real_bundle()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slides = []
    total = 14
    dist = data["dist"]
    ps = data["pressao_state"]
    est = data["est"] or {}
    cui = data["cui"] or {}

    # 1 Capa
    s = prs.slides.add_slide(blank)
    _add_bg(s, GREEN_DARK)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RED
    accent.line.fill.background()
    _textbox(s, 0.7, 1.3, 12, 0.4, "CIEVS-MT · SES-MT · LEITURA AO VIVO DO PAINEL", size=14, color=GREEN_SOFT)
    _textbox(s, 0.7, 1.9, 12, 1.0, "SIS Clima-Saúde MT", size=42, bold=True, color=WHITE)
    _textbox(
        s,
        0.7,
        3.1,
        12,
        1.2,
        "Números reais · Alertas multinível gerados agora\n"
        f"{data['n_mun']} municípios · backend {data['backend']} · {data['gerado_em']}",
        size=18,
        color=RGBColor(0xD7, 0xEB, 0xE3),
    )
    niv_est = str(est.get("nivel") or "—")
    _textbox(
        s,
        0.7,
        5.0,
        12,
        0.8,
        f"Alerta estadual agora: {est.get('icone', '')} {niv_est.upper()} — {est.get('nivel_rotulo', '')}",
        size=20,
        bold=True,
        color=WHITE,
    )
    slides.append(s)

    # 2 Snapshot
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "Snapshot ao vivo do painel", "resumo_municipal_atual + SISREG live")
    _card(s, 0.4, 1.3, 3.0, 2.0, "Municípios", f"{data['n_mun']}\nna base operacional", accent=GREEN)
    _card(s, 3.55, 1.3, 3.0, 2.0, "Tmáx máxima", f"{_fmt_num(data['tmax_max'], 1)} °C\nmédia {_fmt_num(data['tmax_med'], 1)} °C", accent=ORANGE)
    _card(s, 6.7, 1.3, 3.0, 2.0, "UTCI máx.", f"{_fmt_num(data['utci_max'], 1)}\nestresse térmico", accent=RED)
    _card(s, 9.85, 1.3, 3.0, 2.0, "Ocupação méd.", f"{_fmt_num(data['ocup_med'], 1)}%\nIndicaSUS", accent=YELLOW)
    _card(
        s,
        0.4,
        3.55,
        6.2,
        3.0,
        "Distribuição operacional (5 cores)",
        f"🟢 Verde {dist.get('verde', 0)}   🟡 Amarela {dist.get('amarela', 0)}\n"
        f"🟠 Laranja {dist.get('laranja', 0)}   🔴 Vermelha {dist.get('vermelha', 0)}\n"
        f"🟣 Roxa {dist.get('roxa', 0)}   ⚪ Cinza {dist.get('cinza', 0)}\n\n"
        f"Alertas gerados nesta rodada: {data['n_alertas']} pacotes",
        accent=PURPLE,
    )
    _card(
        s,
        6.85,
        3.55,
        6.0,
        3.0,
        "Índice de pressão G/A/V",
        f"Média {_fmt_num(ps.get('indice_media'), 1)} · máx {_fmt_num(ps.get('indice_max'), 1)}\n"
        f"🟢 {ps.get('n_verde', 0)}  🟡 {ps.get('n_amarela', 0)}  🔴 {ps.get('n_vermelha', 0)}\n"
        f"Tendência 7d ↑ {ps.get('n_subindo', 0)} · → {ps.get('n_estavel', 0)} · ↓ {ps.get('n_descendo', 0)}\n"
        f"SISREG: {data['sisreg_n']} mun. · {data['sisreg_fonte']}",
        accent=GREEN,
    )
    _footer(s, 2, total)
    slides.append(s)

    # 3 SISREG
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "SISREG ao vivo na pressão da rede", f"Fonte: {data['sisreg_fonte']}")
    _bullets(
        s,
        0.5,
        1.3,
        12.2,
        1.5,
        [
            f"Municípios com fila/pendência mapeada: {data['sisreg_n']}/142",
            f"Solicitações em aberto (soma estadual): {_fmt_num(data['sisreg_sols_total'], 0)}",
            "Credenciais: SISREG_HOST=10.15.1.71 · DB SES · usuário sisreg_sureg (rede SES)",
        ],
        size=16,
    )
    top_lines = []
    sis_top = data["sisreg_top"]
    if isinstance(sis_top, pd.DataFrame) and not sis_top.empty:
        for _, r in sis_top.iterrows():
            top_lines.append(
                f"{r.get('municipio')}: {_fmt_num(r.get('solicitacoes_abertas'), 0)} abertas · "
                f"espera méd. {_fmt_num(float(r.get('fila_media_h') or 0) / 24.0, 0)} dias"
            )
    _card(s, 0.5, 3.0, 12.3, 3.6, "Maiores filas / pendências (SISREG live)", "\n".join(top_lines[:8]) or "Sem ranking", accent=ORANGE)
    _footer(s, 3, total)
    slides.append(s)

    # 4 Alerta estadual
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    niv = str(est.get("nivel") or "cinza")
    _header_bar(s, "Alerta ① Estadual → SES / CIEVS", est.get("titulo") or "—")
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.25), Inches(12.4), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = LEVEL_RGB.get(niv, MUTED)
    bar.line.fill.background()
    _textbox(s, 0.65, 1.35, 12, 0.5, f"{est.get('icone', '')}  {est.get('nivel_rotulo', niv).upper()}  ·  {est.get('n_municipios', '—')} municípios", size=18, bold=True, color=WHITE)
    _card(s, 0.45, 2.15, 6.1, 2.4, "Motivo / gatilhos", str(est.get("motivo") or "—")[:420], accent=LEVEL_RGB.get(niv, GREEN))
    pred = est.get("predicao") or {}
    _card(s, 6.75, 2.15, 6.1, 2.4, "Predição ~7 dias", str(pred.get("resumo") or "—")[:420], accent=PURPLE)
    _card(s, 0.45, 4.7, 12.4, 2.0, "Indicadores do boletim estadual", "\n".join(_inds_lines(est.get("indicadores"), 7)), accent=GREEN)
    _footer(s, 4, total)
    slides.append(s)

    # 5 Orientações estadual
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "Orientações do alerta estadual", "Gestor · Profissionais · População")
    o = est.get("orientacoes") or {}
    _card(s, 0.4, 1.3, 4.1, 5.2, "Gestor / SES", str(o.get("gestor") or "—"), accent=PURPLE)
    _card(s, 4.65, 1.3, 4.1, 5.2, "Profissionais de saúde", str(o.get("profissional") or "—"), accent=ORANGE)
    _card(s, 8.9, 1.3, 4.0, 5.2, "População", str(o.get("populacao") or "—"), accent=GREEN)
    _footer(s, 5, total)
    slides.append(s)

    # 6 Cuiabá
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    niv = str(cui.get("nivel") or "cinza")
    _header_bar(s, "Alerta ④ Vigidesastre Cuiabá", cui.get("titulo") or "—")
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.25), Inches(12.4), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = LEVEL_RGB.get(niv, MUTED)
    bar.line.fill.background()
    _textbox(s, 0.65, 1.35, 12, 0.5, f"{cui.get('icone', '')}  {cui.get('nivel_rotulo', niv)}  ·  IBGE 5103403", size=18, bold=True, color=WHITE)
    _card(s, 0.45, 2.15, 12.4, 1.6, "Motivo integrado", str(cui.get("motivo") or "—")[:500], accent=LEVEL_RGB.get(niv, RED))
    _card(s, 0.45, 3.95, 6.1, 2.7, "Indicadores Cuiabá", "\n".join(_inds_lines(cui.get("indicadores"), 8)), accent=ORANGE)
    pred = cui.get("predicao") or {}
    oc = cui.get("orientacoes") or {}
    _card(
        s,
        6.75,
        3.95,
        6.1,
        2.7,
        "Predição + orientação ao gestor",
        f"{pred.get('resumo', '—')}\n\nGestor: {str(oc.get('gestor') or '—')[:280]}",
        accent=PURPLE,
    )
    _footer(s, 6, total)
    slides.append(s)

    # 7 Regionais críticas
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "Alerta ② Regionais de Saúde", "Pior nível por regional (gerado agora)")
    lines = []
    for p in data["regs"][:10]:
        lines.append(f"{p.get('icone')} {str(p.get('nivel')).upper()} — {p.get('alvo_nome')} ({p.get('n_municipios')} mun.) · {str(p.get('motivo') or '')[:70]}")
    _bullets(s, 0.5, 1.35, 12.3, 5.4, lines or ["Sem regionais"], size=14)
    _footer(s, 7, total)
    slides.append(s)

    # 8 Municipais críticos
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "Alerta ③ Municípios prioritários", "Pacotes ≥ laranja gerados nesta rodada")
    lines = []
    for p in data["muns"][:12]:
        lines.append(f"{p.get('icone')} {p.get('alvo_nome')} — {p.get('nivel')} · {str(p.get('motivo') or '')[:90]}")
    _bullets(s, 0.5, 1.35, 12.3, 5.4, lines or ["Sem municípios acima do limiar"], size=13)
    _footer(s, 8, total)
    slides.append(s)

    # 9 Exemplo municipal completo (pior)
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    top_m = data["muns"][0] if data["muns"] else {}
    _header_bar(s, "Boletim municipal completo (exemplo)", top_m.get("titulo") or "—")
    _card(s, 0.4, 1.25, 12.5, 1.3, "Motivo", str(top_m.get("motivo") or "—")[:450], accent=LEVEL_RGB.get(str(top_m.get("nivel")), RED))
    _card(s, 0.4, 2.7, 6.2, 3.9, "Indicadores", "\n".join(_inds_lines(top_m.get("indicadores"), 10)), accent=ORANGE)
    om = top_m.get("orientacoes") or {}
    pred = top_m.get("predicao") or {}
    _card(
        s,
        6.8,
        2.7,
        6.1,
        3.9,
        "Predição e orientações",
        f"Predição: {pred.get('resumo', '—')}\n\n"
        f"Gestor: {str(om.get('gestor') or '—')[:200]}\n\n"
        f"Profissional: {str(om.get('profissional') or '—')[:180]}\n\n"
        f"População: {str(om.get('populacao') or '—')[:160]}",
        accent=PURPLE,
    )
    _footer(s, 9, total)
    slides.append(s)

    # 10 Top clima do resumo
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "Municípios com maior score no painel", "Ranking real de resumo_municipal_atual")
    lines = []
    top = data["top_mun"]
    if isinstance(top, pd.DataFrame):
        for _, r in top.iterrows():
            lines.append(
                f"{r.get('municipio')} · {r.get('nivel')} · score {_fmt_num(r.get('score'), 0)} · "
                f"Tmáx {_fmt_num(r.get('tmax'), 1)}°C · UTCI {_fmt_num(r.get('utci_proxy'), 1)} · "
                f"risco3d {_fmt_num(r.get('risco_cumulativo_3d'), 1)}"
            )
    _bullets(s, 0.5, 1.35, 12.3, 5.4, lines[:12], size=13)
    _footer(s, 10, total)
    slides.append(s)

    # 11 Como ler no painel
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "Onde ver isso no painel real", "http://localhost:8501")
    _bullets(
        s,
        0.55,
        1.4,
        12.2,
        5.2,
        [
            "Visão executiva — distribuição de níveis e KPIs estaduais",
            "Mapas — cloropletas municipais (shapefile IBGE)",
            "Assistência — índice de pressão G/A/V + IndicaSUS + SISREG + SINAN + SIM",
            "Alertas — abas ① Estadual · ② Regionais · ③ Municipais · ④ Vigidesastre Cuiabá",
            "Cada boletim: ícone, indicadores, predição 7d, orientações (gestor/profissional/população)",
            "Validar no painel antes de armar SEND_ALERT_ON_LEVEL_CHANGE",
        ],
        size=16,
    )
    _footer(s, 11, total)
    slides.append(s)

    # 12 Fontes
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "Fontes nesta rodada", "Clima + assistência + epidemiologia + regulação")
    _bullets(
        s,
        0.55,
        1.4,
        12.2,
        5.2,
        [
            "Clima: Open-Meteo / TITAN · INMET · Cemaden · ANA · Copernicus/CAMS (quando OK)",
            "Assistência: IndicaSUS (ocupação) · SISREG live (fila/pendências)",
            "Epidemiologia: SINAN · SIVEP · SIM (agravos/óbitos sensíveis ao clima)",
            "Motor de alertas: sisclima/engines/alertas_multinivel.py",
            "Índice de pressão: sisclima/engines/indice_pressao_saude.py",
            f"Gerado em {data['gerado_em']} a partir de {data['backend']}",
        ],
        size=16,
    )
    _footer(s, 12, total)
    slides.append(s)

    # 13 Próximos passos
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "Próximos passos operacionais", "Do insight ao disparo validado")
    _bullets(
        s,
        0.55,
        1.4,
        12.2,
        5.2,
        [
            "Rotina de plantão: abrir painel → Assistência → Alertas (4 níveis)",
            "Manter atualizar_sisreg_pressao.py na carga diária (VPN SES)",
            "Completar CSV de contatos estadual/regional/municipal/Cuiabá",
            "Validar boletins na prévia do painel; só então SEND_ALERT=true",
            "Articular gatilhos do SIS aos níveis do Plano de Contingência seca/estiagem",
        ],
        size=16,
    )
    _footer(s, 13, total)
    slides.append(s)

    # 14 Encerramento
    s = prs.slides.add_slide(blank)
    _add_bg(s, GREEN_DARK)
    _textbox(s, 0.7, 2.0, 12, 0.8, "Decidir com o painel aberto", size=36, bold=True, color=WHITE)
    _textbox(
        s,
        0.7,
        3.1,
        12,
        1.6,
        f"Agora: alerta estadual {est.get('icone', '')} {str(est.get('nivel') or '—').upper()}\n"
        f"Cuiabá Vigidesastre {cui.get('icone', '')} {str(cui.get('nivel') or '—').upper()}\n"
        f"{data['n_alertas']} pacotes prontos para validação.",
        size=18,
        color=RGBColor(0xD7, 0xEB, 0xE3),
    )
    _textbox(s, 0.7, 5.4, 12, 0.8, "CIEVS-MT · SES-MT · http://localhost:8501", size=14, color=RGBColor(0xA8, 0xC7, 0xBC))
    slides.append(s)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))

    # sidecar JSON com alertas reais (para auditoria / anexo)
    side = out.with_suffix(".alertas.json")
    serializable = []
    for p in data["payloads"]:
        serializable.append(
            {
                "escopo": p.get("escopo"),
                "alvo_id": p.get("alvo_id"),
                "alvo_nome": p.get("alvo_nome"),
                "nivel": p.get("nivel"),
                "icone": p.get("icone"),
                "titulo": p.get("titulo"),
                "motivo": p.get("motivo"),
                "predicao": p.get("predicao"),
                "orientacoes": p.get("orientacoes"),
                "indicadores": p.get("indicadores"),
                "n_municipios": p.get("n_municipios"),
                "gerado_em": p.get("gerado_em"),
            }
        )
    side.write_text(
        json.dumps(
            {
                "gerado_em": data["gerado_em"],
                "backend": data["backend"],
                "dist": data["dist"],
                "pressao": data["pressao_state"],
                "sisreg_fonte": data["sisreg_fonte"],
                "n_alertas": data["n_alertas"],
                "alertas": serializable,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    return out, data


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", type=str, default=str(DEFAULT_OUT))
    ap.add_argument("--also-downloads", action="store_true")
    args = ap.parse_args()
    out = Path(args.out)
    path, data = build(out)
    print(f"[OK] {path}")
    print(f"[OK] alertas: {path.with_suffix('.alertas.json')} ({data['n_alertas']} pacotes)")
    print(
        f"[INFO] estadual={data['est'].get('nivel') if data['est'] else '—'} · "
        f"cuiaba={data['cui'].get('nivel') if data['cui'] else '—'} · "
        f"sisreg={data['sisreg_n']} · backend={data['backend']}"
    )
    if args.also_downloads:
        dest = Path.home() / "Downloads" / path.name
        try:
            shutil.copy2(path, dest)
            print(f"[OK] cópia: {dest}")
        except Exception as exc:
            print(f"[WARN] não copiou para Downloads: {exc}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
