# -*- coding: utf-8 -*-
"""
Gera apresentação ~20 min: SIS Clima-Saúde MT (CIEVS/SES-MT).

Contexto: articula o painel operacional ao Plano de Contingência para
Emergência em Saúde Pública por Seca e Estiagem (MT 2026-2027).

Uso:
  .venv\\Scripts\\python.exe gerar_apresentacao_sis_20min.py
  .venv\\Scripts\\python.exe gerar_apresentacao_sis_20min.py --out "C:\\Users\\...\\Desktop\\SIS_20min.pptx"
"""
from __future__ import annotations

import argparse
import shutil
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "docs" / "apresentacoes"
DEFAULT_OUT = OUT_DIR / "Apresentacao_SIS_Clima_Saude_20min.pptx"

# Paleta institucional (painel SIS)
GREEN_DARK = RGBColor(0x08, 0x35, 0x2E)
GREEN = RGBColor(0x0F, 0x6E, 0x56)
GREEN_SOFT = RGBColor(0x1A, 0x8A, 0x6E)
CREAM = RGBColor(0xF3, 0xF6, 0xF4)
INK = RGBColor(0x1A, 0x2E, 0x28)
MUTED = RGBColor(0x4B, 0x5C, 0x56)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
PURPLE = RGBColor(0x6D, 0x28, 0xD9)
YELLOW = RGBColor(0xC4, 0x92, 0x00)

LEVELS = [
    ("Verde", "Normalidade / vigilância de rotina", GREEN_SOFT),
    ("Amarela", "Atenção — desvio inicial", YELLOW),
    ("Laranja", "Alerta — pressão assistencial/climática", ORANGE),
    ("Vermelha", "Alerta alto — resposta ampliada", RED),
    ("Roxa", "Mobilização plena CIEVS / articulação estadual", PURPLE),
]


def _set_run(run, *, size=18, bold=False, color=INK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _add_bg(slide, color: RGBColor):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return shape


def _textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color)
    return box


def _bullets(slide, left, top, width, height, items, *, size=16, color=INK, spacing=True):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if spacing:
            p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"• {item}"
        _set_run(run, size=size, color=color)
    return box


def _notes(slide, text: str):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def _header_bar(slide, title: str, subtitle: str | None = None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN_DARK
    bar.line.fill.background()
    _textbox(slide, 0.45, 0.22, 12.2, 0.45, title, size=26, bold=True, color=WHITE)
    if subtitle:
        _textbox(slide, 0.45, 0.62, 12.2, 0.35, subtitle, size=13, color=RGBColor(0xC8, 0xE6, 0xDC))


def _footer(slide, page: int, total: int):
    _textbox(
        slide,
        0.45,
        7.1,
        10,
        0.3,
        "CIEVS-MT · SIS Clima-Saúde · Código legível · Fontes oficiais",
        size=10,
        color=MUTED,
    )
    _textbox(slide, 11.2, 7.1, 1.7, 0.3, f"{page}/{total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def _card(slide, left, top, width, height, title, body, accent=GREEN):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = RGBColor(0xD5, 0xE0, 0xDA)
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.12), Inches(height)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()
    _textbox(slide, left + 0.28, top + 0.15, width - 0.4, 0.35, title, size=14, bold=True, color=GREEN_DARK)
    _textbox(slide, left + 0.28, top + 0.5, width - 0.4, height - 0.65, body, size=13, color=INK)


def load_live_snapshot() -> dict:
    """Métricas opcionais do banco; se falhar, usa placeholders honestos."""
    snap = {
        "backend": "—",
        "n_mun": 142,
        "nivel_estado": "—",
        "mun_critico": "—",
        "dist": {},
        "ok": False,
    }
    try:
        from sisclima.core.db import backend_name, read_table

        resumo = read_table("resumo_municipal_atual")
        if resumo is None or resumo.empty:
            return snap
        snap["backend"] = backend_name()
        snap["n_mun"] = int(resumo["cod_ibge"].nunique()) if "cod_ibge" in resumo.columns else len(resumo)
        if "nivel" in resumo.columns:
            niv = resumo["nivel"].astype(str).str.lower()
            snap["dist"] = {k: int((niv == k).sum()) for k, _, __ in LEVELS for k in [k.lower()]}
            # sentinel: maior score se existir
            if "score" in resumo.columns:
                top = resumo.sort_values("score", ascending=False).iloc[0]
                snap["nivel_estado"] = str(top.get("nivel", "—"))
                snap["mun_critico"] = str(top.get("municipio", "—"))
            else:
                # pior nível presente
                order = ["roxa", "vermelha", "laranja", "amarela", "verde"]
                for o in order:
                    if (niv == o).any():
                        snap["nivel_estado"] = o
                        break
        snap["ok"] = True
    except Exception as exc:
        snap["error"] = str(exc)
    return snap


def build(out: Path) -> Path:
    snap = load_live_snapshot()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 14
    slides_meta = []

    # 1 Capa
    s = prs.slides.add_slide(blank)
    _add_bg(s, GREEN_DARK)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN_SOFT
    accent.line.fill.background()
    _textbox(s, 0.7, 1.6, 11.5, 0.4, "CIEVS-MT · SES-MT", size=16, color=GREEN_SOFT)
    _textbox(s, 0.7, 2.2, 11.5, 1.2, "SIS Clima-Saúde MT", size=44, bold=True, color=WHITE)
    _textbox(
        s,
        0.7,
        3.5,
        11.5,
        1.0,
        "Vigilância integrada clima–saúde para apoiar o plantão,\no Plano de Contingência (seca e estiagem) e a resposta estadual.",
        size=20,
        color=RGBColor(0xD7, 0xEB, 0xE3),
    )
    _textbox(
        s,
        0.7,
        5.5,
        11.5,
        0.6,
        f"Apresentação operacional · ~20 minutos · {date.today().strftime('%d/%m/%Y')}",
        size=14,
        color=RGBColor(0xA8, 0xC7, 0xBC),
    )
    _notes(
        s,
        "ROTEIRO (1 min): Apresente-se (CIEVS). Diga que o SIS é a ferramenta de monitoramento "
        "que alimenta a leitura de risco clima–saúde alinhada ao Plano de Contingência seca/estiagem 2026-2027.",
    )
    slides_meta.append(s)

    # 2 Agenda + tempo
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "Roteiro (~20 minutos)", "Uma pergunta por bloco")
    _bullets(
        s,
        0.6,
        1.4,
        12,
        5.2,
        [
            "Por que clima e saúde juntos em Mato Grosso? (3 min)",
            "O que é o SIS Clima-Saúde e o que ele NÃO é (3 min)",
            "Fontes, níveis operacionais e mapas municipais (5 min)",
            "TITAN, solo, alertas e AdaptaSUS no plantão (5 min)",
            "Governança técnica, próximos passos e perguntas (4 min)",
        ],
        size=20,
    )
    _footer(s, 2, total)
    _notes(s, "ROTEIRO (0,5–1 min): Mostre o relógio mental. Prometa demonstração conceitual do painel, não tour longo de telas.")
    slides_meta.append(s)

    # 3 Contexto plano
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "Contexto estadual", "Plano de Contingência — seca, estiagem e calor (2026–2027)")
    _bullets(
        s,
        0.55,
        1.35,
        7.4,
        5.2,
        [
            "Objetivo do Plano: coordenar vigilância e atenção para mitigar impactos de seca/estiagem.",
            "Base: CME/CIEVS, Sala de Situação SES-MT, decreto de emergência por incêndios (2024).",
            "Cenário: seca extrema, calor >40 °C, focos de calor e sobrecarga assistencial no período seco.",
            "Níveis de resposta do Plano (0 → IV) exigem indicadores contínuos e territoriais.",
        ],
        size=17,
    )
    _card(
        s,
        8.2,
        1.5,
        4.6,
        4.6,
        "Papel do SIS",
        "Traduz clima, ar, água, solo, assistência e agravos em um nível operacional municipal "
        "para o CIEVS decidir com a mesma lógica de preparação/resposta do Plano.",
        accent=GREEN,
    )
    _footer(s, 3, total)
    _notes(
        s,
        "ROTEIRO (3 min): Use o PPT do Plano (DOC-WA) como âncora política. Não releia o plano — "
        "diga que o SIS é o radar diário que sustenta aqueles eixos e níveis.",
    )
    slides_meta.append(s)

    # 4 Problema
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "O problema operacional", "Dados existem — a decisão precisa ser integrada e oportuna")
    _card(s, 0.5, 1.4, 4.0, 4.8, "Fragmentação", "Clima (INMET/Open-Meteo), queimadas/ar, Cemaden, ANA, IndicaSUS/CNES, SIVEP, SINAN e sentinela vivem em silos diferentes.", accent=ORANGE)
    _card(s, 4.7, 1.4, 4.0, 4.8, "Tempo do plantão", "O CIEVS precisa de leitura estadual + municipal em minutos, com motivo, orientação e fila de atenção — não de dezenas de planilhas.", accent=YELLOW)
    _card(s, 8.9, 1.4, 4.0, 4.8, "Resposta escalonada", "Os níveis do Plano (mobilização → crise) pedem evidência rastreável para acionar vigilâncias, assistência e comunicação.", accent=RED)
    _footer(s, 4, total)
    _notes(s, "ROTEIRO (2 min): Conte um exemplo real de plantão (calor + fumaça + ocupação).")
    slides_meta.append(s)

    # 5 O que é
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "O que é o SIS Clima-Saúde", "Painel único de vigilância integrada — CIEVS-MT")
    _bullets(
        s,
        0.55,
        1.35,
        12.2,
        5.0,
        [
            "Sistema operacional que cruza exposição climática/ambiental com desfechos e capacidade assistencial.",
            "Cobre os 142 municípios de MT com mapas por polígono (shapefile), filtros por regional e município.",
            "Produz nível operacional (verde→roxa), score, motivos, orientação em linguagem simples e predição ~7 dias.",
            "Incorpora camada TITAN (calor, solo, INMET/Cemaden) e alinhamento AdaptaSUS / guia MS.",
            "Não substitui boletins oficiais de longo prazo (ex.: cenário de setembro) — complementa com horizonte curto.",
        ],
        size=17,
    )
    _footer(s, 5, total)
    _notes(s, "ROTEIRO (3 min): Repita o limite: predição ~7 dias ≠ forecast sazonal oficial.")
    slides_meta.append(s)

    # 6 Fontes
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "Fontes e arquitetura", "APIs/SQL oficiais · código legível · auditável na rede SES")
    sources = [
        ("Clima", "Open-Meteo (Tmax, umidade, vento, precipitação, solo)"),
        ("Alertas", "INMET · Cemaden (wsAlertas2) · ANA telemetria"),
        ("Ar", "Copernicus/CAMS (quando disponível) + proxies"),
        ("Saúde", "DW SINAN/SIVEP · IndicaSUS/CNES · sentinela SG"),
        ("Território", "Shapefile IBGE/MT · população · regionais de saúde"),
        ("Entrega", "PostgreSQL/Docker · Streamlit · alertas e-mail/Telegram"),
    ]
    y = 1.35
    for i, (t, b) in enumerate(sources):
        col = i % 3
        row = i // 3
        _card(s, 0.45 + col * 4.25, y + row * 2.5, 4.05, 2.2, t, b, accent=GREEN if i % 2 == 0 else GREEN_SOFT)
    _footer(s, 6, total)
    _notes(
        s,
        "ROTEIRO (2,5 min): Enfatize política SES: sem scrapers ofuscados; User-Agent SIS-Clima-Saude-MT. "
        "DW quando a rede permitir; fallback CSV/sample documentado.",
    )
    slides_meta.append(s)

    # 7 Níveis
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "Níveis operacionais do SIS", "Leitura compatível com a lógica de resposta do Plano")
    y = 1.35
    for name, desc, color in LEVELS:
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(y), Inches(12.2), Inches(0.95))
        bar.fill.solid()
        bar.fill.fore_color.rgb = WHITE
        bar.line.color.rgb = RGBColor(0xD5, 0xE0, 0xDA)
        chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y + 0.2), Inches(2.1), Inches(0.55))
        chip.fill.solid()
        chip.fill.fore_color.rgb = color
        chip.line.fill.background()
        _textbox(s, 0.75, y + 0.28, 2.0, 0.4, name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _textbox(s, 3.1, y + 0.28, 9.3, 0.5, desc, size=16, color=INK)
        y += 1.05
    _footer(s, 7, total)
    _notes(
        s,
        "ROTEIRO (2 min): Relacione Verde≈Nível 0; Amarela/Laranja≈Mobilização/Alerta; "
        "Vermelha/Roxa≈Emergência/necessidade de articulação ampliada — sem afirmar equivalência jurídica automática.",
    )
    slides_meta.append(s)

    # 8 Snapshot ao vivo
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "Situação na rodada atual", "Snapshot do banco operacional (se disponível)")
    if snap["ok"]:
        dist_txt = ", ".join(f"{k}: {v}" for k, v in snap["dist"].items() if v)
        body = (
            f"Base: {snap['backend']} · Municípios: {snap['n_mun']}\n"
            f"Município mais crítico (score): {snap['mun_critico']}\n"
            f"Nível no topo: {snap['nivel_estado']}\n"
            f"Distribuição: {dist_txt or '—'}"
        )
    else:
        body = (
            "Abra o painel local (Docker/Streamlit) para exibir o snapshot ao vivo.\n"
            "Em rodadas recentes: 142 municípios com mapas por shapefile e alerta integrado."
        )
    _card(s, 0.55, 1.5, 12.2, 2.4, "Leitura executiva", body, accent=PURPLE)
    _bullets(
        s,
        0.7,
        4.2,
        12,
        2.4,
        [
            "No painel: Visão executiva → Mapas → Clima/TITAN → Alertas.",
            "Filtros por Regional de Saúde e Município no topo (sem menu lateral confuso).",
            "Cada cor no mapa = polígono municipal (não pontos).",
        ],
        size=16,
    )
    _footer(s, 8, total)
    _notes(s, "ROTEIRO (2 min): Se possível, alterne para o navegador (localhost:8501) por 60–90 s.")
    slides_meta.append(s)

    # 9 Mapas
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "Mapas cloropléticos", "Shapefile municipal MT — decisão territorial")
    _bullets(
        s,
        0.55,
        1.4,
        12.2,
        5.0,
        [
            "Geometria oficial dos municípios (shapefile / GeoJSON processado).",
            "Indicadores: nível, tensão climática, carga em saúde, ocupação, PM2,5, solo, arboviroses, SRAG.",
            "Evita leitura por ‘pontinhos’ — mostra área real do município para regionais e gestores.",
            "Recorte por regional preserva o contexto estadual no topo do painel.",
        ],
        size=18,
    )
    _footer(s, 9, total)
    _notes(s, "ROTEIRO (2 min): Mostre um mapa de nível e um de tensão climática.")
    slides_meta.append(s)

    # 10 TITAN / alertas
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "TITAN + alerta integrado", "Calor, solo, hidrologia e canais oficiais no mesmo máximo")
    _card(s, 0.5, 1.4, 6.0, 4.9, "Camada TITAN", "UTCI/proxy, ondas de calor, saturação do solo (Open-Meteo), resiliência de leitos (IndicaSUS+CNES), risco hidro (ANA).", accent=GREEN)
    _card(s, 6.8, 1.4, 6.0, 4.9, "Alerta integrado", "nivel_alerta_integrado = máximo entre SIS, INMET, Cemaden, solo, hidro e calor — com justificativa em linguagem de plantão.", accent=ORANGE)
    _footer(s, 10, total)
    _notes(s, "ROTEIRO (3 min): Explique o ‘máximo’ como regra conservadora de segurança sanitária.")
    slides_meta.append(s)

    # 11 AdaptaSUS / MS
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "AdaptaSUS e indicadores MS", "Alinhamento conceitual — sem inventar protocolo paralelo")
    _bullets(
        s,
        0.55,
        1.4,
        12.2,
        5.0,
        [
            "Aba AdaptaSUS / Guia MS: riscos climáticos e recomendações de adaptação em saúde.",
            "SIVEP e Sentinela SG com catálogo alinhado a indicadores MS/SVSA.",
            "Sazonalidade e odds ratio para apoiar leitura epidemiológica–climática (com cautela causal).",
            "Ajudante de interpretação: cards de guia + texto justificativo (padrão Meningites/CIEVS).",
        ],
        size=17,
    )
    _footer(s, 11, total)
    _notes(s, "ROTEIRO (2 min): Deixe claro: o SIS apoia; a decisão clínica/gestora continua humana.")
    slides_meta.append(s)

    # 12 Governança
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "Governança técnica na rede SES", "Transparência para TI e para a vigilância")
    _bullets(
        s,
        0.55,
        1.4,
        12.2,
        5.0,
        [
            "Proibido: código ofuscado/esfumaçado, scrapers stealth, bypass de WAF/CAPTCHA.",
            "Obrigatório: Python legível, SQL versionado em sql/, User-Agent institucional identificável.",
            "Segredos só em .env / secrets do Streamlit — nunca no GitHub.",
            "Deploy: GitHub (branch painel-v9) · Streamlit Community Cloud · Docker local SES.",
        ],
        size=17,
    )
    _footer(s, 12, total)
    _notes(s, "ROTEIRO (1,5 min): Mensagem para TI: tráfego auditável reduz risco de bloqueio.")
    slides_meta.append(s)

    # 13 Próximos passos
    s = prs.slides.add_slide(blank)
    _add_bg(s, CREAM)
    _header_bar(s, "Próximos passos", "Do painel à rotina do CIEVS")
    _bullets(
        s,
        0.55,
        1.4,
        12.2,
        5.0,
        [
            "Padronizar abertura do painel no plantão (checklist de 5 minutos).",
            "Integrar envio de alertas (e-mail/Telegram) com validação humana.",
            "Treinar regionais na leitura de mapas e níveis (linguagem simples).",
            "Manter atualização do enriquecimento operacional e da base Postgres.",
            "Articular indicadores do SIS aos gatilhos do Plano de Contingência (mesa técnica).",
        ],
        size=18,
    )
    _footer(s, 13, total)
    _notes(s, "ROTEIRO (2 min): Peça donos: CIEVS (rotina), TI (rede), assistência (ocupação).")
    slides_meta.append(s)

    # 14 Encerramento
    s = prs.slides.add_slide(blank)
    _add_bg(s, GREEN_DARK)
    _textbox(s, 0.7, 2.0, 12, 0.8, "Obrigado", size=40, bold=True, color=WHITE)
    _textbox(
        s,
        0.7,
        3.0,
        12,
        1.5,
        "SIS Clima-Saúde MT — vigilância para decidir a tempo.\n"
        "Perguntas e encaminhamentos.",
        size=20,
        color=RGBColor(0xD7, 0xEB, 0xE3),
    )
    _textbox(
        s,
        0.7,
        5.2,
        12,
        1.0,
        "CIEVS-MT · SES-MT\nPainel: streamlit_app.py · Documentação em docs/",
        size=14,
        color=RGBColor(0xA8, 0xC7, 0xBC),
    )
    _notes(s, "ROTEIRO (2–3 min): Abra para perguntas. Ofereça demonstração guiada após a reunião.")
    slides_meta.append(s)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def main() -> int:
    ap = argparse.ArgumentParser(description="Gera apresentação SIS ~20 min")
    ap.add_argument("--out", type=str, default=str(DEFAULT_OUT))
    ap.add_argument("--also-downloads", action="store_true", help="Copia também para Downloads")
    args = ap.parse_args()
    out = Path(args.out)
    path = build(out)
    print(f"[OK] {path}")
    if args.also_downloads:
        dest = Path.home() / "Downloads" / path.name
        shutil.copy2(path, dest)
        print(f"[OK] cópia: {dest}")
    # roteiro texto
    roteiro = OUT_DIR / "ROTEIRO_APRESENTACAO_SIS_20min.md"
    roteiro.write_text(
        """# Roteiro — SIS Clima-Saúde (~20 min)

| Min | Slide | Fala-chave |
|-----|-------|------------|
| 0–1 | Capa | Quem somos; o SIS apoia o plantão e o Plano de Contingência |
| 1–2 | Roteiro | Cinco blocos; sem tour infinito de telas |
| 2–5 | Contexto | Plano seca/estiagem 2026–2027 exige indicadores contínuos |
| 5–7 | Problema | Silos de dados × tempo do plantão × resposta escalonada |
| 7–10 | O que é | 142 municípios; nível; 7 dias ≠ sazonal |
| 10–12 | Fontes | Oficiais + código legível na rede SES |
| 12–14 | Níveis | Verde→roxa e relação prudente com níveis do Plano |
| 14–16 | Snapshot/Mapas | Demo rápida do painel + cloropletas |
| 16–19 | TITAN/AdaptaSUS | Máximo integrado + guia MS |
| 19–20 | Governança/próximos | TI, rotina CIEVS, perguntas |

**Dica:** tenha o PPT do Plano (DOC-WA) e o painel `http://localhost:8501` abertos em abas.
""",
        encoding="utf-8",
    )
    print(f"[OK] {roteiro}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
