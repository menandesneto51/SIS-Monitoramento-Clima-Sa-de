# -*- coding: utf-8 -*-
"""
Gera apresentação CIEVS — cenário operacional setembro/2026 a partir do ARARAS.

Limite técnico explícito: predição numérica ~7 dias ≠ forecast climático mensal.
O cenário de setembro é leitura sazonal/prospectiva apoiada nos indicadores atuais.

Template: Slide Padrão.pptx | Citações ABNT NBR 6023:2018 quando houver fontes oficiais.

Uso:
  .venv\\Scripts\\python.exe gerar_apresentacao_cenario_setembro.py
  .venv\\Scripts\\python.exe gerar_apresentacao_cenario_setembro.py --template "C:\\...\\Slide Padrão.pptx"
"""
from __future__ import annotations

import argparse
import shutil
from datetime import date
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

from sisclima.core.db import backend_name, read_table

ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "docs" / "apresentacoes"
ASSETS = OUT_DIR / "assets"
DEFAULT_OUT = OUT_DIR / "Apresentacao_SIS_Cenario_Setembro_2026.pptx"
ACCESS = date(2026, 7, 27)

LEVEL_ORDER = ["verde", "amarela", "laranja", "vermelha", "roxa", "cinza"]
LEVEL_COLORS = {
    "verde": "#16803c",
    "amarela": "#c49200",
    "laranja": "#d97706",
    "vermelha": "#dc2626",
    "roxa": "#6d28d9",
    "cinza": "#6b7280",
}

CITE = {
    "painel": "(INMET et al., 2026a)",
    "inpe_not": "(INPE, 2026)",
    "sesmt": "(MATO GROSSO, [2024?])",
    "sis": "(CIEVS-MT, 2026)",
    "om": "(OPEN-METEO, 2026)",
}


def find_template(explicit: str | None = None) -> Path:
    if explicit:
        p = Path(explicit)
        if p.exists():
            return p
        raise FileNotFoundError(f"Template não encontrado: {explicit}")
    downloads = Path(r"C:\Users\Menandesneto\Downloads")
    candidates = [
        downloads / "Slide Padrão.pptx",
        downloads / "Slide Padrao.pptx",
    ]
    if downloads.exists():
        candidates.extend(sorted(downloads.glob("Slide Padr*.pptx")))
        candidates.extend(sorted(downloads.glob("*Padrao*.pptx")))
    for c in candidates:
        if c.exists():
            return c
    raise FileNotFoundError("Template 'Slide Padrão.pptx' não encontrado. Use --template.")


def _safe_numeric(s) -> pd.Series:
    return pd.to_numeric(s, errors="coerce")


def _lvl_counts(series: pd.Series) -> pd.Series:
    return series.astype(str).str.lower().value_counts().reindex(LEVEL_ORDER).fillna(0).astype(int)


def load_metrics() -> dict:
    """Extrai métricas reais do Postgres ARARAS para a apresentação."""
    resumo = read_table("resumo_municipal_atual")
    pred = read_table("predicao_calor_7d_municipal_v6")
    corr = read_table("analise_clima_saude_correlacoes_v8")
    alerta = read_table("alerta_inteligente_municipal_v6")
    cemaden = read_table("cemaden_alertas")
    ana = read_table("ana_risco_municipal")
    try:
        pipeline = read_table("pipeline_runs")
    except Exception:
        pipeline = pd.DataFrame()

    if resumo.empty:
        raise RuntimeError("resumo_municipal_atual vazio — rode o pipeline/enriquecimento antes.")

    if "cod_ibge" in resumo.columns:
        resumo = resumo.drop_duplicates("cod_ibge", keep="first")

    niveis = _lvl_counts(resumo["nivel"]) if "nivel" in resumo.columns else pd.Series(dtype=int)

    sort_cols = [c for c in ["score", "risco_cumulativo_3d", "pressao_calor_pct"] if c in resumo.columns]
    top = resumo.copy()
    if sort_cols:
        for c in sort_cols:
            top[c] = _safe_numeric(top[c])
        top = top.sort_values(sort_cols, ascending=[False] * len(sort_cols)).head(10)
    sent = top.iloc[0] if not top.empty else None

    pred_dist = pd.Series(dtype=int)
    pred_top = pd.DataFrame()
    if not pred.empty and "nivel_predicao_7d" in pred.columns:
        pred_dist = _lvl_counts(pred["nivel_predicao_7d"])
        if "risco_preditivo_score" in pred.columns:
            pred = pred.copy()
            pred["risco_preditivo_score"] = _safe_numeric(pred["risco_preditivo_score"])
            pred_top = pred.sort_values("risco_preditivo_score", ascending=False).head(10)

    corr_top = pd.DataFrame()
    if not corr.empty and "abs_rho" in corr.columns:
        corr = corr.copy()
        corr["abs_rho"] = _safe_numeric(corr["abs_rho"])
        corr_top = corr.sort_values("abs_rho", ascending=False).head(10)

    alerta_dist = pd.Series(dtype=int)
    recom_dist = pd.Series(dtype=int)
    if not alerta.empty:
        if "alerta_inteligente_nivel" in alerta.columns:
            alerta_dist = _lvl_counts(alerta["alerta_inteligente_nivel"])
        if "recomendacao_operacional" in alerta.columns:
            recom_dist = alerta["recomendacao_operacional"].astype(str).value_counts().head(6)

    srag_top = pd.DataFrame()
    if "casos_srag" in resumo.columns:
        tmp = resumo.copy()
        tmp["casos_srag"] = _safe_numeric(tmp["casos_srag"])
        srag_top = tmp.sort_values("casos_srag", ascending=False).head(8)

    ar_top = pd.DataFrame()
    if "pm25_ugm3" in resumo.columns:
        tmp = resumo.copy()
        tmp["pm25_ugm3"] = _safe_numeric(tmp["pm25_ugm3"])
        ar_top = tmp.dropna(subset=["pm25_ugm3"]).sort_values("pm25_ugm3", ascending=False).head(8)

    coverage = {}
    for c in ["tmax", "casos_srag", "pm25_ugm3", "casos_arbovirus_7d", "ocupacao_leitos_pct", "pressao_calor_pct"]:
        coverage[c] = int(_safe_numeric(resumo[c]).notna().sum()) if c in resumo.columns else 0

    onda_n = 0
    if "onda_calor_p95_2d" in resumo.columns:
        onda_n = int((_safe_numeric(resumo["onda_calor_p95_2d"]).fillna(0) > 0).sum())

    return {
        "resumo": resumo,
        "pred": pred,
        "corr": corr,
        "alerta": alerta,
        "cemaden": cemaden,
        "ana": ana,
        "pipeline": pipeline,
        "niveis": niveis,
        "top": top,
        "sent": sent,
        "pred_dist": pred_dist,
        "pred_top": pred_top,
        "corr_top": corr_top,
        "alerta_dist": alerta_dist,
        "recom_dist": recom_dist,
        "srag_top": srag_top,
        "ar_top": ar_top,
        "coverage": coverage,
        "onda_n": onda_n,
        "backend": backend_name(),
    }


def make_charts(m: dict) -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    paths: dict[str, Path] = {}

    niveis = m["niveis"]
    if not niveis.empty:
        fig, ax = plt.subplots(figsize=(8, 4.2))
        labels = [x.capitalize() for x in LEVEL_ORDER if int(niveis.get(x, 0)) > 0]
        vals = [int(niveis.get(x.lower(), 0)) for x in labels]
        colors = [LEVEL_COLORS.get(lb.lower(), "#334155") for lb in labels]
        ax.bar(labels, vals, color=colors)
        ax.set_title("ARARAS — municípios por nível operacional")
        ax.set_ylabel("Municípios")
        for i, v in enumerate(vals):
            ax.text(i, v + 0.5, str(v), ha="center", va="bottom", fontsize=9)
        fig.tight_layout()
        p = ASSETS / "distribuicao_niveis.png"
        fig.savefig(p, dpi=160)
        plt.close(fig)
        paths["niveis"] = p

    top = m["top"]
    if not top.empty and "municipio" in top.columns and "score" in top.columns:
        plot = top.head(10).iloc[::-1]
        fig, ax = plt.subplots(figsize=(8.5, 5))
        ax.barh(plot["municipio"].astype(str), _safe_numeric(plot["score"]).fillna(0), color="#0b3d34")
        ax.set_xlabel("Score operacional")
        ax.set_title("ARARAS — ranking municípios críticos (score)")
        fig.tight_layout()
        p = ASSETS / "ranking_score.png"
        fig.savefig(p, dpi=160)
        plt.close(fig)
        paths["ranking"] = p

    pred_dist = m["pred_dist"]
    if not pred_dist.empty and pred_dist.sum() > 0:
        fig, ax = plt.subplots(figsize=(8, 4.2))
        labels = [x.capitalize() for x in LEVEL_ORDER if int(pred_dist.get(x, 0)) > 0]
        vals = [int(pred_dist.get(x.lower(), 0)) for x in labels]
        colors = [LEVEL_COLORS.get(lb.lower(), "#334155") for lb in labels]
        ax.bar(labels, vals, color=colors)
        ax.set_title("ARARAS — predição operacional ~7 dias (não sazonal)")
        ax.set_ylabel("Municípios")
        fig.tight_layout()
        p = ASSETS / "predicao_7d_dist.png"
        fig.savefig(p, dpi=160)
        plt.close(fig)
        paths["pred"] = p

    corr_top = m["corr_top"]
    if not corr_top.empty and {"exposicao", "desfecho", "abs_rho"}.issubset(corr_top.columns):
        plot = corr_top.head(8).copy()
        plot["par"] = plot["exposicao"].astype(str) + " → " + plot["desfecho"].astype(str)
        plot = plot.sort_values("abs_rho", ascending=True)
        fig, ax = plt.subplots(figsize=(9, 4.8))
        ax.barh(plot["par"], plot["abs_rho"], color="#0f6e56")
        ax.set_xlabel("|ρ| Spearman (exploratório)")
        ax.set_title("ARARAS — correlação clima–saúde (não causal)")
        fig.tight_layout()
        p = ASSETS / "correlacoes_top.png"
        fig.savefig(p, dpi=160)
        plt.close(fig)
        paths["corr"] = p

    srag_top = m["srag_top"]
    if not srag_top.empty and "municipio" in srag_top.columns:
        plot = srag_top.head(8).iloc[::-1]
        fig, ax = plt.subplots(figsize=(8.5, 4.6))
        ax.barh(plot["municipio"].astype(str), _safe_numeric(plot["casos_srag"]).fillna(0), color="#9f1239")
        ax.set_xlabel("Casos SRAG")
        ax.set_title("ARARAS — municípios com mais casos SRAG")
        fig.tight_layout()
        p = ASSETS / "srag_top.png"
        fig.savefig(p, dpi=160)
        plt.close(fig)
        paths["srag"] = p

    return paths


def _clear_slides(prs: Presentation) -> None:
    while len(prs.slides) > 0:
        sldId = prs.slides._sldIdLst[0]
        prs.part.drop_rel(sldId.rId)
        del prs.slides._sldIdLst[0]


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = title
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            shape.text = subtitle
            break


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str], font_pt: int = 15) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(font_pt)
        p.font.name = "Calibri"


def add_two_content_slide(
    prs: Presentation,
    title: str,
    left_bullets: list[str],
    right_image: Path | None = None,
    right_bullets: list[str] | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    slide.shapes.title.text = title
    left = slide.placeholders[1].text_frame
    left.clear()
    for i, line in enumerate(left_bullets):
        p = left.paragraphs[0] if i == 0 else left.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(13)
    right_ph = slide.placeholders[2]
    if right_image and right_image.exists():
        sp = right_ph.element
        sp.getparent().remove(sp)
        slide.shapes.add_picture(str(right_image), Inches(5.1), Inches(1.55), width=Inches(4.5))
    elif right_bullets:
        tf = right_ph.text_frame
        tf.clear()
        for i, line in enumerate(right_bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(12)


def build_presentation(template: Path, out_path: Path) -> Path:
    """Monta PPTX conforme estrutura do plano (cenário setembro)."""
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    m = load_metrics()
    charts = make_charts(m)

    shutil.copy2(template, out_path)
    prs = Presentation(str(out_path))
    _clear_slides(prs)

    resumo = m["resumo"]
    niveis = m["niveis"]
    sent = m["sent"]
    cov = m["coverage"]
    mun_n = int(resumo["cod_ibge"].nunique()) if "cod_ibge" in resumo.columns else len(resumo)
    tmax = _safe_numeric(resumo["tmax"]).max() if "tmax" in resumo.columns else np.nan
    umid = _safe_numeric(resumo["umidade_media"]).mean() if "umidade_media" in resumo.columns else np.nan
    precip = _safe_numeric(resumo["precipitacao_mm"]).mean() if "precipitacao_mm" in resumo.columns else np.nan
    risco3 = _safe_numeric(resumo["risco_cumulativo_3d"]).mean() if "risco_cumulativo_3d" in resumo.columns else np.nan
    srag_sum = _safe_numeric(resumo["casos_srag"]).sum() if "casos_srag" in resumo.columns else 0.0
    arbo_sum = _safe_numeric(resumo["casos_arbovirus_7d"]).sum() if "casos_arbovirus_7d" in resumo.columns else 0.0
    sent_mun = str(sent.get("municipio", "—")) if sent is not None else "—"
    sent_nivel = str(sent.get("nivel", "—")).upper() if sent is not None else "—"
    nivel_txt = ", ".join(f"{k.capitalize()} {int(v)}" for k, v in niveis.items() if int(v) > 0)

    pipe_txt = "—"
    pipe = m["pipeline"]
    if isinstance(pipe, pd.DataFrame) and not pipe.empty and "started_at" in pipe.columns:
        last = pipe.sort_values("started_at", ascending=False).iloc[0]
        pipe_txt = f"{last.get('started_at', '—')} · {last.get('status', '—')} · {last.get('message', '—')}"

    # --- 1 Capa ---
    add_title_slide(
        prs,
        "ARARAS MT",
        "Cenário operacional setembro 2026\n"
        "CIEVS-MT · Vigilância integrada clima–saúde\n"
        f"Extração {ACCESS.strftime('%d/%m/%Y')} · {mun_n} municípios",
    )

    # --- 2 Objetivo + disclaimer ---
    add_bullets_slide(
        prs,
        "Objetivo e disclaimer metodológico",
        [
            "Objetivo: apresentar o cenário operacional prospectivo para set/2026 com o que o ARARAS já produz.",
            "DISCLAIMER: a predição numérica do ARARAS cobre cerca de 7 dias — não é forecast climático mensal. "
            + CITE["sis"]
            + " "
            + CITE["om"],
            "O ‘cenário de setembro’ = (1) achados atuais do ARARAS + (2) riscos sazonais típicos de setembro no MT "
            "+ (3) o que o painel passará a vigiar no mês.",
            "Associações clima–saúde são ecológicas/exploratórias — não causalidade individual. " + CITE["sis"],
            "Roteiro de fala: ‘7 dias = número do sistema; setembro = leitura sazonal operacional’.",
        ],
    )

    # --- 3 O que o ARARAS integra ---
    add_two_content_slide(
        prs,
        "O que o ARARAS integra",
        [
            "Clima/biometeo (Open-Meteo): Tmax, umidade, UTCI, risco cumulativo, onda de calor P95.",
            "Saúde: SIVEP/SRAG, arboviroses, Sentinela SG, GeoCalor.",
            "Ambiente: qualidade do ar (PM2,5/PM10/O3 quando disponível).",
            "Hidrologia: Cemaden e ANA.",
            "Inteligência: score/nível Verde→Roxa, correlação Spearman, predição 7d, alerta inteligente.",
            f"Backend: {m['backend']} · última pipeline: {pipe_txt}. " + CITE["sis"],
        ],
        right_bullets=[
            "Painel Streamlit — blocos:",
            "Situação · Clima · Saúde",
            "Hidrologia · Operação · Método",
            "Filtros: regional e município",
            "Mapas por shapefile IBGE",
            f"Cobertura clima: {cov.get('tmax', 0)}/{mun_n}",
            f"SRAG: {cov.get('casos_srag', 0)}/{mun_n}",
            f"PM2,5: {cov.get('pm25_ugm3', 0)}/{mun_n}",
            f"Ocupação leitos: {cov.get('ocupacao_leitos_pct', 0)}/{mun_n}",
        ],
    )

    # --- 4 Situação atual ---
    add_two_content_slide(
        prs,
        "Situação atual do Estado (baseline ARARAS)",
        [
            f"Municípios monitorados: {mun_n}.",
            f"Distribuição por nível: {nivel_txt}. " + CITE["sis"],
            f"Sentinela / mais crítico nesta rodada: {sent_mun} ({sent_nivel}).",
            f"Tmax máx.: {tmax:.1f} °C · UR média: {umid:.0f}% · chuva média: {precip:.2f} mm"
            if pd.notna(tmax)
            else "Clima: —",
            f"Risco cumulativo 3d médio: {risco3:.1f} · municípios em onda P95≥2d: {m['onda_n']}."
            if pd.notna(risco3)
            else f"Onda de calor: {m['onda_n']} mun.",
            f"SRAG (soma municipal): ~{srag_sum:.0f} casos · arbovírus 7d: ~{arbo_sum:.0f} onde há dado.",
            "Uso: priorizar atenção — não declarar epidemia/desastre só pelo score.",
        ],
        right_image=charts.get("niveis"),
    )

    # --- 5 Ranking críticos ---
    top_lines = [
        "Municípios com maior score nesta extração (priorização CIEVS): " + CITE["sis"],
    ]
    for _, row in m["top"].head(8).iterrows():
        top_lines.append(
            f"{row.get('municipio', '—')} — {str(row.get('nivel', '—')).upper()} · "
            f"score {row.get('score', '—')} · risco3d {row.get('risco_cumulativo_3d', '—')}"
        )
    top_lines.append("Mapa detalhado: painel Situação / Mapas (shapefile municipal).")
    add_two_content_slide(
        prs,
        "Municípios mais críticos (ranking ARARAS)",
        top_lines,
        right_image=charts.get("ranking"),
    )

    # --- 6 Por que setembro no MT ---
    add_bullets_slide(
        prs,
        "Por que setembro no MT (sazonalidade operacional)",
        [
            "Setembro no Centro-Oeste: final da seca, calor intenso, pico de queimadas e piora da qualidade do ar.",
            "Contexto oficial JAS/2026: maior probabilidade de chuva abaixo da média no centro-norte e "
            "temperatura acima da normal — potencial de ondas de calor e queimadas. "
            + CITE["painel"]
            + " "
            + CITE["inpe_not"],
            "Indicadores atuais do ARARAS já apontam pressão térmica e SRAG concentrada em polos urbanos. "
            + CITE["sis"],
            "Queimadas × internações respiratórias: leitura de preparação estadual, não previsão de casos. "
            + CITE["sesmt"],
            "Arboviroses: calor e irregularidade hídrica favoráveis ao vetor — monitorar via painel, sem extrapolar epidemia só do clima.",
            "O ARARAS não ‘prevê setembro’: ele organiza vigilância contínua para o mês.",
        ],
        font_pt=14,
    )

    # --- 7 Predição 7d ---
    pred_lines = [
        "Projeção NUMÉRICA disponível agora: predicao_calor_7d_municipal_v6 (~7 dias). " + CITE["sis"],
        "Útil para a semana seguinte — não para ‘projetar setembro inteiro’.",
    ]
    if not m["pred_dist"].empty:
        pred_lines.append(
            "Distribuição prevista: "
            + ", ".join(f"{k.capitalize()} {int(v)}" for k, v in m["pred_dist"].items() if int(v) > 0)
        )
    for _, row in m["pred_top"].head(6).iterrows():
        pred_lines.append(
            f"{row.get('municipio')}: {str(row.get('nivel_predicao_7d', '—')).upper()} · "
            f"Tmax7d {row.get('tmax_max_7d', '—')} °C"
        )
    add_two_content_slide(
        prs,
        "Projeções disponíveis agora (predição ~7 dias)",
        pred_lines,
        right_image=charts.get("pred"),
    )

    # --- 8 Correlação ---
    corr_lines = [
        "Tabela analise_clima_saude_correlacoes_v8 — Spearman municipal. " + CITE["sis"],
        f"Pares calculados nesta base: {len(m['corr'])}.",
        "Interpretação: associação ecológica; confundidores e defasagens não resolvidos aqui.",
    ]
    for _, row in m["corr_top"].head(6).iterrows():
        corr_lines.append(
            f"{row.get('exposicao')} → {row.get('desfecho')}: |ρ|={float(row.get('abs_rho', 0)):.2f} "
            f"(n={row.get('n_municipios', '—')})"
        )
    add_two_content_slide(
        prs,
        "Correlação clima–saúde (exploratória)",
        corr_lines,
        right_image=charts.get("corr"),
    )

    # --- 9 Eixos de risco set/2026 ---
    eixos = [
        "Eixos que o painel deve vigiar em set/2026 (leitura operacional): " + CITE["sis"],
        f"1) Calor — risco3d médio {risco3:.1f}; onda P95: {m['onda_n']} mun.; acompanhar UTCI e ondas."
        if pd.notna(risco3)
        else f"1) Calor — onda P95: {m['onda_n']} mun.",
        f"2) Qualidade do ar — PM2,5 em {cov.get('pm25_ugm3', 0)} mun. (cobertura parcial).",
        f"3) Arboviroses — ~{arbo_sum:.0f} casos/7d onde há dado ({cov.get('casos_arbovirus_7d', 0)} mun.).",
        f"4) SRAG — ~{srag_sum:.0f} casos na base; polos: "
        + ", ".join(str(x) for x in m["srag_top"]["municipio"].head(4).tolist())
        if not m["srag_top"].empty
        else f"4) SRAG — ~{srag_sum:.0f} casos na base.",
        f"5) Hidrologia — Cemaden {len(m['cemaden'])} reg.; ANA {len(m['ana'])} mun.",
        "Contexto sazonal (oficial): estiagem/queimadas no JAS — preparar resposta, não inventar magnitude. "
        + CITE["painel"],
    ]
    add_two_content_slide(
        prs,
        "Eixos de risco para set/2026",
        eixos,
        right_image=charts.get("srag") or charts.get("niveis"),
    )

    # --- 10 Recomendações ---
    rec_lines = [
        "Saídas do módulo alerta_inteligente_municipal_v6: " + CITE["sis"],
    ]
    if not m["alerta_dist"].empty:
        rec_lines.append(
            "Alerta inteligente: "
            + ", ".join(f"{k.capitalize()} {int(v)}" for k, v in m["alerta_dist"].items() if int(v) > 0)
        )
    for rec, n in m["recom_dist"].items():
        rec_lines.append(f"{n} municípios → {rec}")
    rec_lines.extend(
        [
            "CIEVS/regionais: usar ranking + alerta para sala de situação e plantão.",
            "Comunicar calor, ar e SRAG com base nos dados do dia — atualizar após cada pipeline.",
        ]
    )
    add_bullets_slide(prs, "Recomendações operacionais CIEVS / regionais", rec_lines, font_pt=13)

    # --- 11 Lacunas e próximos passos ---
    add_bullets_slide(
        prs,
        "Lacunas e próximos passos",
        [
            f"IndicaSUS / ocupação de leitos: {cov.get('ocupacao_leitos_pct', 0)}/{mun_n} nesta rodada "
            "(indisponível sem credenciais válidas) — painel usa proxy clima+saúde com ressalva.",
            f"Qualidade do ar: cobertura parcial ({cov.get('pm25_ugm3', 0)} mun.) — não generalizar ao estado todo.",
            "Não há modelo climático mensal (ECMWF/seasonal) no ARARAS — fora de escopo desta entrega.",
            "Próximos passos: validar IndicaSUS; manter pipeline diário; regenerar slides após atualizar o banco.",
            "Comando: python gerar_apresentacao_cenario_setembro.py",
            "Referências oficiais (El Niño/JAS): docs/apresentacoes/REFERENCIAS_ABNT_6023.md",
        ],
        font_pt=14,
    )

    # --- 12 Encerramento ---
    add_title_slide(
        prs,
        "Obrigado",
        "CIEVS-MT · ARARAS MT\n"
        "Cenário operacional setembro 2026\n"
        "Painel Streamlit · predicao 7d ≠ forecast mensal\n"
        f"Extração {ACCESS.strftime('%d/%m/%Y')}",
    )

    prs.save(str(out_path))
    return out_path


def main() -> int:
    parser = argparse.ArgumentParser(description="Gera PPTX cenário setembro a partir do ARARAS")
    parser.add_argument("--template", default=None, help="Caminho do Slide Padrão.pptx")
    parser.add_argument("--out", default=str(DEFAULT_OUT), help="Arquivo PPTX de saída")
    args = parser.parse_args()
    template = find_template(args.template)
    out = Path(args.out)
    print(f"[INFO] Template: {template}")
    print(f"[INFO] Saída: {out}")
    print(f"[INFO] Backend: {backend_name()}")
    path = build_presentation(template, out)
    n_slides = len(Presentation(str(path)).slides)
    print(f"[OK] {path}")
    print(f"[OK] Slides: {n_slides}")
    print(f"[OK] Assets: {ASSETS}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
