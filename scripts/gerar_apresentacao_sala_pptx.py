"""Gera PPTX editável da Sala de Situação (SE) a partir do resumo ARARAS + e-SUS + STAR."""
from __future__ import annotations

import argparse
import json
from pathlib import Path

import pandas as pd

ROOT = Path(__file__).resolve().parents[1]


def _fmt(n, nd=0, suf=""):
    if n is None or (isinstance(n, float) and pd.isna(n)):
        return "—"
    try:
        v = float(n)
    except (TypeError, ValueError):
        return str(n)
    if nd == 0:
        return f"{int(round(v)):,}".replace(",", ".") + suf
    return f"{v:.{nd}f}".replace(".", ",") + suf


def _add_title(slide, text: str, subtitle: str = ""):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(1.0))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x13, 0x51, 0xB4)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def _bullets(slide, lines: list[str], top: float = 1.5):
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.0), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        p.level = 0
        p.font.size = Pt(16)
        p.space_after = Pt(8)


def build_pptx(out: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    from sisclima.core.db import read_table
    from sisclima.engines.boletim_el_nino.snapshot import merge_predicao_7d, snapshot_operacional
    from sisclima.engines.esus_clima_analise import analisar_esus_clima

    resumo = read_table("resumo_municipal_atual")
    try:
        pred = read_table("predicao_calor_7d_municipal_v6")
    except Exception:
        pred = None
    snap = snapshot_operacional(merge_predicao_7d(resumo, pred))
    esus = analisar_esus_clima()
    star = {}
    star_path = ROOT / "data" / "output" / "star" / "STAR_geocalor_carga.json"
    if star_path.exists():
        star = json.loads(star_path.read_text(encoding="utf-8"))
    star_resumo = {}
    sr = ROOT / "data" / "output" / "star" / "STAR_resumo_indicadores.json"
    if sr.exists():
        star_resumo = json.loads(sr.read_text(encoding="utf-8"))

    n = snap.get("n_municipios")
    vr = snap.get("n_vermelha_roxa")
    proj = snap.get("niveis_projecao_7d") or {}
    vr_proj = int(proj.get("vermelha") or 0) + int(proj.get("roxa") or 0)
    ext = (snap.get("extremos") or {}).get("tmax") or {}

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 capa
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(
        s,
        "Sala de Situação CIEVS-MT — ARARAS MT",
        "Boletim El Niño · SE 35/2026 · material editável",
    )
    _bullets(
        s,
        [
            "Portaria n.º 0590/2026/GBSES",
            "Fonte operacional: ARARAS MT / CIEVS-MT",
            "Complementos: e-SUS APS (Centralizador) · STAR/GeoCalor (CDS ERA5-Land)",
            "Associação ecológica — não implica causalidade individual",
        ],
        top=2.0,
    )

    # 2 cartões
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s, "Cartões executivos", "Rodada operacional")
    _bullets(
        s,
        [
            f"Vermelho/roxo atual: {_fmt(vr)}/{_fmt(n)}",
            f"Projeção ~7 dias: {_fmt(vr_proj)}/{_fmt(n)}",
            f"Tmáx máxima: {_fmt(ext.get('tmax'), 1, ' °C')}",
            f"Mun. Tmáx ≥ 37 °C: {_fmt(snap.get('n_tmax_37'))}",
            f"Mun. PM2,5 ≥ 25: {_fmt(snap.get('n_pm25_25'))}",
            f"Focos 7d (referência): {_fmt(snap.get('focos_7d_total'))}",
        ],
    )

    # 3 e-SUS
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s, "Atenção primária — e-SUS APS × clima", "Centralizador PEC/eSUS · cruzamento ARARAS")
    if esus.get("ok"):
        corrs = esus.get("correlacoes") or {}
        c1 = corrs.get("atend_28d_x_tmax") or {}
        c2 = corrs.get("resp_cid_28d_x_pm25") or {}
        _bullets(
            s,
            [
                f"Municípios: {_fmt(esus.get('n'))} · críticos VR: {_fmt(esus.get('n_criticos'))}",
                f"Cadastro asma {_fmt(esus.get('asma'))} · idosos 60+ {_fmt(esus.get('idoso_60mais'))} · gestantes {_fmt(esus.get('gestante'))}",
                f"Atendimentos 28d {_fmt(esus.get('atendimentos_28d'))} · CID resp. 28d {_fmt(esus.get('resp_cid_28d'))}",
                f"Média asma (VR vs demais): {_fmt(esus.get('media_asma_criticos'), 0)} vs {_fmt(esus.get('media_asma_outros'), 0)}",
                f"Spearman atend.28d × Tmáx: ρ={_fmt(c1.get('rho'), 2)} (n={_fmt(c1.get('n'))})",
                f"Spearman CID resp. × PM2,5: ρ={_fmt(c2.get('rho'), 2)} (n={_fmt(c2.get('n'))})",
            ],
        )
    else:
        _bullets(s, ["Dados e-SUS indisponíveis nesta geração."])

    # 4 STAR
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s, "Ondas de calor — STAR / GeoCalor", "EHF Nairn & Fawcett · CDS ERA5-Land")
    _bullets(
        s,
        [
            f"Flag operacional onda (P95≥2d): {_fmt(snap.get('n_onda_calor') or star_resumo.get('n_onda_flag'))} municípios",
            f"Tmáx máx. (STAR/resumo): {_fmt(star_resumo.get('tmax_max') or ext.get('tmax'), 1, ' °C')}",
            f"Catálogo EHF: {star.get('metodologia') or 'GeoCalor_EHF_NairnFawcett_3d'}",
            f"Janela carga: {star.get('janela') or '—'}",
            f"Eventos EHF: {_fmt(star.get('n_eventos'))} · dias de onda: {_fmt(star.get('n_dias_onda'))}",
            f"Fonte: {star.get('fonte_clima') or 'Copernicus CDS / cache local'}",
            "GeoCalor público não cobre Cuiabá/MT — ARARAS aplica EHF aos 142 municípios",
        ],
    )

    # 5 prioridades
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s, "Prioridades para a Sala", "Encaminhamentos operacionais")
    _bullets(
        s,
        [
            "Manter vigilância de calor/desidratação e agravos respiratórios em VR",
            "APS: busca ativa em idosos, asma/DPOC, gestantes e acamados nos críticos",
            "Cruzar ocupação IndicaSUS com pressão térmica e fumaça (não inventar % estadual)",
            "Acompanhar carga CDS 2020–2026 e atualizar Anexo STAR quando concluir",
            "Reavaliar projeção ~7d na próxima rodada (chuva recente × retorno do calor)",
        ],
    )

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument(
        "--out",
        type=Path,
        default=ROOT / "docs" / "apresentacoes" / "Sala_Situacao_SE35_2026_ARARAS.pptx",
    )
    args = ap.parse_args()
    path = build_pptx(args.out)
    print(path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
