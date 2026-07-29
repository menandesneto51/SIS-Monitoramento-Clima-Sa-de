# -*- coding: utf-8 -*-
"""
Apresentação panorâmica SIS Clima-Saúde MT
- Alinhada ao Plano de Contingência (DOC-20260611-WA0013)
- Dados reais do Postgres + prints do painel
- Prévia do alerta Telegram/e-mail (envio OFF)

Uso:
  .venv\\Scripts\\python.exe gerar_apresentacao_sis_panorama.py
  .venv\\Scripts\\python.exe gerar_apresentacao_sis_panorama.py --also-downloads
"""
from __future__ import annotations

import argparse
import json
import shutil
from datetime import datetime
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "docs" / "apresentacoes"
PRINTS = OUT_DIR / "prints_painel"
DEFAULT_OUT = OUT_DIR / "Apresentacao_SIS_Panorama_Indicadores_Alertas.pptx"

GREEN_DARK = RGBColor(0x08, 0x35, 0x2E)
GREEN = RGBColor(0x0F, 0x6E, 0x56)
GREEN_SOFT = RGBColor(0x1A, 0x8A, 0x6E)
CREAM = RGBColor(0xF3, 0xF6, 0xF4)
INK = RGBColor(0x1A, 0x2E, 0x28)
MUTED = RGBColor(0x4B, 0x5C, 0x56)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
PURPLE = RGBColor(0x6D, 0x28, 0xD9)
YELLOW = RGBColor(0xC4, 0x92, 0x00)

NAV_SECTIONS = [
    "Visão executiva", "Mapas", "Guia do leitor", "Clima / TITAN", "Qualidade do ar",
    "Assistência", "Arboviroses", "SIVEP", "Sentinela SG", "GeoCalor",
    "AdaptaSUS / Guia MS", "Correlação clima-saúde", "Cemaden / ANA", "Sazonalidade / OR",
    "Operacional", "Geografia", "Inteligência", "Alertas", "Cálculos",
]

LEVEL_RGB = {
    "verde": GREEN_SOFT, "amarela": YELLOW, "laranja": ORANGE,
    "vermelha": RED, "roxa": PURPLE, "cinza": MUTED,
}


def _run(run, *, size=16, bold=False, color=INK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _bg(slide, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(13.333), Inches(7.5))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    spTree = slide.shapes._spTree
    el = sh._element
    spTree.remove(el)
    spTree.insert(2, el)


def _txt(slide, l, t, w, h, text, *, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _run(r, size=size, bold=bold, color=color)
    return box


def _bullets(slide, l, t, w, h, items, *, size=14, color=INK):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        r = p.add_run()
        r.text = f"• {item}"
        _run(r, size=size, color=color)
    return box


def _header(slide, title, subtitle=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN_DARK
    bar.line.fill.background()
    _txt(slide, 0.45, 0.15, 12.2, 0.4, title, size=22, bold=True, color=WHITE)
    if subtitle:
        _txt(slide, 0.45, 0.55, 12.2, 0.35, subtitle, size=12, color=RGBColor(0xD7, 0xEB, 0xE3))


def _footer(slide, n, total):
    _txt(slide, 0.45, 7.15, 10, 0.25, "SIS Clima-Saúde MT · CIEVS/SES-MT · dados reais Postgres", size=9, color=MUTED)
    _txt(slide, 11.2, 7.15, 1.7, 0.25, f"{n}/{total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def _card(slide, l, t, w, h, title, body, accent=GREEN):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = RGBColor(0xD5, 0xE0, 0xDA)
    ab = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(0.1), Inches(h))
    ab.fill.solid()
    ab.fill.fore_color.rgb = accent
    ab.line.fill.background()
    _txt(slide, l + 0.22, t + 0.1, w - 0.35, 0.3, title, size=12, bold=True, color=GREEN_DARK)
    _txt(slide, l + 0.22, t + 0.42, w - 0.35, h - 0.55, body, size=11, color=INK)


def _add_pic(slide, path: Path, l, t, w, h):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))
        return True
    _txt(slide, l, t, w, 0.4, f"[Print ausente: {path.name}]", size=11, color=RED)
    return False


def _fmt(v, nd=1):
    try:
        if v is None or (isinstance(v, float) and pd.isna(v)):
            return "—"
        x = float(v)
        if abs(x) >= 1000:
            return f"{x:,.0f}".replace(",", ".")
        return f"{x:.{nd}f}"
    except Exception:
        return str(v)


def load_bundle() -> dict:
    from sisclima.core.config import env
    from sisclima.core.db import backend_name, read_table
    from sisclima.engines.alertas_multinivel import (
        build_alertas_multinivel,
        render_payload_markdown,
    )
    from sisclima.engines.indice_pressao_saude import state_pressao_summary
    from sisclima.alerts.change_detector import build_level_change_message

    resumo = read_table("resumo_municipal_atual")
    pred = read_table("predicao_calor_7d_municipal_v6")
    ai = read_table("alerta_integrado_sis_titan")
    sis = read_table("ops_sisreg_municipio")
    press = read_table("indice_pressao_saude_municipal_v1")
    hist = read_table("alertas_enviados")

    if resumo is None or resumo.empty:
        raise RuntimeError("resumo_municipal_atual vazio")

    payloads = build_alertas_multinivel(
        resumo,
        alerta_integrado=ai if not ai.empty else None,
        predicao_7d=pred if not pred.empty else None,
        min_level="laranja",
    )
    est = next(p for p in payloads if p.get("escopo") == "estadual")
    cui = next(p for p in payloads if p.get("escopo") == "cuiaba")
    muns = sorted(
        [p for p in payloads if p.get("escopo") == "municipal"],
        key=lambda p: {"roxa": 4, "vermelha": 3, "laranja": 2}.get(str(p.get("nivel")), 0),
        reverse=True,
    )[:10]

    niv = resumo["nivel"].astype(str).str.lower().value_counts().to_dict()
    press_state = state_pressao_summary(press) if not press.empty else {}

    # Prévia Telegram/e-mail: multinível (completo) + formato clássico mudança de nível
    md_est = render_payload_markdown(est)
    md_cui = render_payload_markdown(cui)
    # Telegram limita ~4096 chars — cortamos com marcador
    tg_body = md_est[:3800] + ("\n\n…[mensagem truncada para Telegram]" if len(md_est) > 3800 else "")
    subj_legacy, msg_legacy = build_level_change_message(
        datetime.now().date().isoformat(),
        "laranja",
        str(est.get("nivel")),
        str(est.get("motivo") or "").split("; "),
    )

    send_flag = env("SEND_ALERT_ON_LEVEL_CHANGE", "false")
    email_on = env("ALERT_EMAIL_ENABLED", "false")
    tg_on = env("ALERT_TELEGRAM_ENABLED", "false")
    real_send = False
    hist_note = "Nenhum registro em alertas_enviados (envio nunca disparado nesta base)."
    if hist is not None and not hist.empty:
        hist_note = f"{len(hist)} registro(s) em alertas_enviados."
        if "status" in hist.columns and hist["status"].astype(str).str.contains("enviado", case=False, na=False).any():
            real_send = True

    return {
        "gerado_em": datetime.now().isoformat(timespec="seconds"),
        "backend": backend_name(),
        "dist": niv,
        "tmax_max": float(pd.to_numeric(resumo.get("tmax"), errors="coerce").max()),
        "utci_max": float(pd.to_numeric(resumo.get("utci_proxy"), errors="coerce").max()) if "utci_proxy" in resumo.columns else None,
        "ocup_med": float(pd.to_numeric(resumo.get("ocupacao_leitos_pct"), errors="coerce").mean()) if "ocupacao_leitos_pct" in resumo.columns else None,
        "tensao": float(pd.to_numeric(resumo.get("indice_tensao_climatica"), errors="coerce").mean()) if "indice_tensao_climatica" in resumo.columns else None,
        "carga": float(pd.to_numeric(resumo.get("indice_carga_saude"), errors="coerce").mean()) if "indice_carga_saude" in resumo.columns else None,
        "vig": float(pd.to_numeric(resumo.get("indice_vigilancia_integrada"), errors="coerce").mean()) if "indice_vigilancia_integrada" in resumo.columns else None,
        "sisreg_n": int(len(sis)),
        "sisreg_sols": float(pd.to_numeric(sis.get("solicitacoes_abertas"), errors="coerce").sum()) if not sis.empty else None,
        "pressao": press_state,
        "est": est,
        "cui": cui,
        "muns": muns,
        "n_alertas": len(payloads),
        "md_est": md_est,
        "md_cui": md_cui,
        "tg_body": tg_body,
        "email_subject": f"[SIS Clima-Saúde / CIEVS] {est.get('titulo')}",
        "email_body": md_est,
        "legacy_subject": subj_legacy,
        "legacy_body": msg_legacy,
        "send_flag": send_flag,
        "email_on": email_on,
        "tg_on": tg_on,
        "real_send": real_send,
        "hist_note": hist_note,
        "canais_preview": "Telegram + E-mail (credenciais presentes; SEND_ALERT_ON_LEVEL_CHANGE=false)",
    }


def build(out: Path) -> tuple[Path, dict]:
    data = load_bundle()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 14
    n = 0
    dist = data["dist"]
    est = data["est"]
    cui = data["cui"]

    def new():
        nonlocal n
        n += 1
        s = prs.slides.add_slide(blank)
        return s

    # 1 Capa
    s = new()
    _bg(s, GREEN_DARK)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PURPLE
    accent.line.fill.background()
    _txt(s, 0.6, 1.4, 12, 0.35, "CIEVS-MT · SES-MT · Plano de Contingência seca/estiagem 2026-2027", size=13, color=GREEN_SOFT)
    _txt(s, 0.6, 1.9, 12, 0.9, "SIS Clima-Saúde MT", size=40, bold=True, color=WHITE)
    _txt(
        s, 0.6, 3.0, 12, 1.2,
        "Panorama do projeto · Indicadores · Abas do painel · Alertas reais\n"
        f"Rodada {data['gerado_em']} · backend {data['backend']} · 142 municípios",
        size=16, color=RGBColor(0xD7, 0xEB, 0xE3),
    )
    _txt(
        s, 0.6, 5.0, 12, 0.8,
        f"Alerta estadual agora: {est.get('icone')} {str(est.get('nivel')).upper()}  ·  "
        f"Cuiabá: {cui.get('icone')} {str(cui.get('nivel')).upper()}",
        size=18, bold=True, color=WHITE,
    )
    _footer(s, n, total)

    # 2 Panorama
    s = new()
    _bg(s, CREAM)
    _header(s, "Panorama do projeto", "O que é o SIS e para que serve na decisão")
    _bullets(
        s, 0.5, 1.25, 12.3, 5.5,
        [
            "Painel único CIEVS/SES-MT para vigilância integrada clima–saúde em 142 municípios.",
            "Une fontes oficiais: Open-Meteo/TITAN, INMET, Cemaden, ANA, SINAN, SIVEP, SIM, IndicaSUS, SISREG.",
            "Produz níveis operacionais (Verde→Roxa), índice de pressão G/A/V e predição ~7 dias.",
            "Gera alertas em 4 níveis: estadual (SES), regional, municipal e Vigidesastre Cuiabá.",
            "Apoia o Plano de Contingência seca/estiagem — o SIS informa e prioriza; a gestão decide COE/portarias.",
            "Código legível (sem ofuscação); painel em http://localhost:8501 (Docker Postgres).",
            f"Nesta rodada: Roxa {dist.get('roxa',0)} · Vermelha {dist.get('vermelha',0)} · "
            f"Laranja {dist.get('laranja',0)} · Amarela {dist.get('amarela',0)} · Verde {dist.get('verde',0)}.",
        ],
        size=15,
    )
    _footer(s, n, total)

    # 3 Plano × SIS
    s = new()
    _bg(s, CREAM)
    _header(s, "Tabela Plano × SIS (template DOC-20260611-WA0013)", "Níveis Zero→IV ↔ Verde→Roxa")
    _card(s, 0.4, 1.2, 12.5, 5.5, "Ponte operacional",
          "Zero — Normalidade  ↔  Verde / Amarela baixa  → rotina\n"
          "I — Mobilização  ↔  Amarela / Laranja  → boletim e articulação regional\n"
          "II — Alerta  ↔  Laranja / Vermelha  → prioridade operacional + pressão assistencial\n"
          "III — Emergência  ↔  Vermelha / Roxa  → sala de situação / COE\n"
          "IV — Crise  ↔  Roxa sustentada + múltiplos gatilhos  → mobilização plena\n\n"
          f"Leitura desta rodada: estadual {str(est.get('nivel')).upper()} e Cuiabá {str(cui.get('nivel')).upper()} "
          "→ alinhamento prudente aos níveis III–IV do Plano (validar gatilhos humanos).",
          accent=PURPLE)
    _footer(s, n, total)

    # 4 Snapshot KPIs
    s = new()
    _bg(s, CREAM)
    _header(s, "Snapshot ao vivo (Postgres)", "Indicadores da rodada atual")
    _card(s, 0.4, 1.2, 3.0, 1.7, "Tmáx máx.", f"{_fmt(data['tmax_max'])} °C", accent=ORANGE)
    _card(s, 3.55, 1.2, 3.0, 1.7, "UTCI máx.", f"{_fmt(data['utci_max'])}", accent=RED)
    _card(s, 6.7, 1.2, 3.0, 1.7, "Ocupação méd.", f"{_fmt(data['ocup_med'])}%", accent=YELLOW)
    _card(s, 9.85, 1.2, 3.0, 1.7, "SISREG", f"{data['sisreg_n']} mun.\n{_fmt(data['sisreg_sols'],0)} abertas", accent=GREEN)
    _card(s, 0.4, 3.15, 4.1, 3.4, "Índices compostos (média estadual)",
          f"Tensão climática: {_fmt(data['tensao'],0)}\n"
          f"Carga em saúde: {_fmt(data['carga'],0)}\n"
          f"Vigilância integrada: {_fmt(data['vig'],0)}\n"
          f"Pressão G/A/V: 🟢 {data['pressao'].get('n_verde',0)} · "
          f"🟡 {data['pressao'].get('n_amarela',0)} · 🔴 {data['pressao'].get('n_vermelha',0)}\n"
          f"Tendência ↑ {data['pressao'].get('n_subindo',0)} · → {data['pressao'].get('n_estavel',0)} · ↓ {data['pressao'].get('n_descendo',0)}",
          accent=GREEN)
    _card(s, 4.7, 3.15, 4.1, 3.4, "Distribuição operacional",
          f"🟢 Verde {dist.get('verde',0)}\n🟡 Amarela {dist.get('amarela',0)}\n"
          f"🟠 Laranja {dist.get('laranja',0)}\n🔴 Vermelha {dist.get('vermelha',0)}\n"
          f"🟣 Roxa {dist.get('roxa',0)}\n\nAlertas gerados: {data['n_alertas']} pacotes",
          accent=PURPLE)
    _card(s, 8.95, 3.15, 3.9, 3.4, "Alerta sentinela",
          f"{est.get('icone')} Estadual {str(est.get('nivel')).upper()}\n"
          f"{cui.get('icone')} Cuiabá {str(cui.get('nivel')).upper()}\n\n"
          f"Motivo estadual:\n{str(est.get('motivo') or '—')[:220]}",
          accent=RED)
    _footer(s, n, total)

    # 5 Indicadores criados
    s = new()
    _bg(s, CREAM)
    _header(s, "Indicadores já criados no SIS", "Catálogo operacional do painel")
    _bullets(
        s, 0.5, 1.25, 12.3, 5.5,
        [
            "Clima/TITAN: Tmáx, UTCI/proxy, risco cumulativo 3d, onda de calor P95, saturação de solo, precipitação.",
            "Ar: PM2,5 / PM10 / O3 / IQA (cobertura parcial Copernicus/CAMS).",
            "Assistência: ocupação IndicaSUS, pressão assistencial, leitos CNES, resiliência.",
            "Índice de pressão G/A/V: IndicaSUS + SISREG + SINAN + SIM (atual, pred 7d, tendência ↑/→/↓).",
            "Epidemiologia: arboviroses 7d/incidência/z-score; SRAG/SIVEP; óbitos SIM calor/cardiorrespiratório.",
            "Compostos: índice_tensao_climatica, indice_carga_saude, indice_vigilancia_integrada, AdaptaSUS.",
            "Estatística: correlação clima–saúde, Odds Ratio, sazonalidade, lags; predição ~7d municipal/regional.",
            "Alertas: nível SIS, alerta integrado SIS+TITAN, alertas multinível (4 escopos).",
        ],
        size=14,
    )
    _footer(s, n, total)

    # 6 Abas do painel
    s = new()
    _bg(s, CREAM)
    _header(s, "Abas do painel (app_v9)", f"{len(NAV_SECTIONS)} seções · navegação horizontal")
    mid = (len(NAV_SECTIONS) + 1) // 2
    _bullets(s, 0.5, 1.25, 6.0, 5.5, NAV_SECTIONS[:mid], size=14)
    _bullets(s, 6.8, 1.25, 6.0, 5.5, NAV_SECTIONS[mid:], size=14)
    _footer(s, n, total)

    # 7 Print Visão
    s = new()
    _bg(s, CREAM)
    _header(s, "Print real — Visão executiva", "Painel http://localhost:8501")
    pic = PRINTS / "A_visao_executiva.png"
    if not pic.exists():
        pic = PRINTS / "01_visao_executiva.png"
    _add_pic(s, pic, 0.4, 1.15, 12.5, 5.7)
    _footer(s, n, total)

    # 8 Print Mapas
    s = new()
    _bg(s, CREAM)
    _header(s, "Print real — Mapas / territorial", "Cloropletas por shapefile municipal")
    pic = PRINTS / "B_mapas.png"
    if not pic.exists():
        pic = PRINTS / "A_visao_mapa.png"
    _add_pic(s, pic, 0.4, 1.15, 12.5, 5.7)
    _footer(s, n, total)

    # 9 Print Assistência
    s = new()
    _bg(s, CREAM)
    _header(s, "Print real — Assistência / índice de pressão", "IndicaSUS · SISREG · SINAN · SIM")
    pic = PRINTS / "C_assistencia_pressao.png"
    if not pic.exists():
        pic = PRINTS / "04_assistencia_conteudo.png"
    _add_pic(s, pic, 0.4, 1.15, 12.5, 5.7)
    _footer(s, n, total)

    # 10 Print Alertas
    s = new()
    _bg(s, CREAM)
    _header(s, "Print real — Aba Alertas", "Prévia multinível no painel")
    pic = PRINTS / "D_alertas_estadual.png"
    if not pic.exists() or pic.stat().st_size < 50_000:
        pic = PRINTS / "D_alertas.png"
    _add_pic(s, pic, 0.4, 1.15, 12.5, 5.7)
    _footer(s, n, total)

    # 11 Alerta estadual completo
    s = new()
    _bg(s, CREAM)
    niv = str(est.get("nivel") or "cinza")
    _header(s, "Modelo de alerta real — ① Estadual (SES/CIEVS)", est.get("titulo") or "—")
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = LEVEL_RGB.get(niv, MUTED)
    bar.line.fill.background()
    _txt(s, 0.55, 1.22, 12.2, 0.4, f"{est.get('icone','')}  {est.get('nivel_rotulo', niv)}  ·  {est.get('n_municipios')} municípios", size=15, bold=True, color=WHITE)
    inds = est.get("indicadores") or []
    ind_txt = "\n".join(f"{i.get('rotulo')}: {i.get('valor')}" for i in inds[:8])
    _card(s, 0.4, 1.9, 6.2, 2.4, "Motivo / gatilhos", str(est.get("motivo") or "—")[:450], accent=LEVEL_RGB.get(niv, RED))
    pred = est.get("predicao") or {}
    _card(s, 6.8, 1.9, 6.1, 2.4, "Predição ~7 dias", str(pred.get("resumo") or "—")[:450], accent=PURPLE)
    _card(s, 0.4, 4.5, 6.2, 2.3, "Indicadores do boletim", ind_txt or "—", accent=GREEN)
    o = est.get("orientacoes") or {}
    _card(s, 6.8, 4.5, 6.1, 2.3, "Orientações",
          f"Gestor: {str(o.get('gestor') or '—')[:160]}\n\n"
          f"Profissional: {str(o.get('profissional') or '—')[:140]}\n\n"
          f"População: {str(o.get('populacao') or '—')[:120]}",
          accent=ORANGE)
    _footer(s, n, total)

    # 12 Cuiabá + top municipais
    s = new()
    _bg(s, CREAM)
    _header(s, "④ Vigidesastre Cuiabá + municípios prioritários", "Pacotes gerados nesta rodada")
    _card(s, 0.4, 1.2, 12.5, 1.8, cui.get("titulo") or "Cuiabá",
          f"Motivo: {str(cui.get('motivo') or '—')[:400]}\n"
          f"Predição: {(cui.get('predicao') or {}).get('resumo', '—')}",
          accent=LEVEL_RGB.get(str(cui.get("nivel")), PURPLE))
    lines = [f"{p.get('icone')} {p.get('alvo_nome')} — {p.get('nivel')} · {str(p.get('motivo') or '')[:90]}" for p in data["muns"][:8]]
    _bullets(s, 0.5, 3.3, 12.3, 3.5, lines, size=12)
    _footer(s, n, total)

    # 13 Prévia Telegram / e-mail
    s = new()
    _bg(s, CREAM)
    _header(s, "Modelo de mensagem Telegram / E-mail", "PRÉVIA DO CANAL — envio OFF (não disparado)")
    status = (
        f"SEND_ALERT_ON_LEVEL_CHANGE={data['send_flag']} · "
        f"ALERT_EMAIL_ENABLED={data['email_on']} · ALERT_TELEGRAM_ENABLED={data['tg_on']}\n"
        f"{data['hist_note']} · {data['canais_preview']}"
    )
    _card(s, 0.4, 1.15, 12.5, 1.15, "Status dos canais", status, accent=ORANGE)
    _card(s, 0.4, 2.5, 12.5, 0.7, "Assunto (e-mail)", data["email_subject"][:200], accent=GREEN)
    # corpo prévia
    preview = data["tg_body"].replace("# ", "").replace("**", "")
    _card(s, 0.4, 3.35, 12.5, 3.4, "Corpo (prévia Telegram ≤~4k caracteres)", preview[:1100] + ("…" if len(preview) > 1100 else ""), accent=PURPLE)
    _footer(s, n, total)

    # 14 Encerramento
    s = new()
    _bg(s, GREEN_DARK)
    _txt(s, 0.6, 2.0, 12, 0.7, "Decidir com o painel aberto", size=32, bold=True, color=WHITE)
    _txt(
        s, 0.6, 3.0, 12, 2.0,
        f"Estadual {est.get('icone')} {str(est.get('nivel')).upper()} · "
        f"Cuiabá {cui.get('icone')} {str(cui.get('nivel')).upper()}\n"
        f"{data['n_alertas']} pacotes prontos · validar antes de armar SEND_ALERT\n"
        "http://localhost:8501 · CIEVS-MT / SES-MT",
        size=16, color=RGBColor(0xD7, 0xEB, 0xE3),
    )
    _footer(s, n, total)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))

    side = {
        "gerado_em": data["gerado_em"],
        "backend": data["backend"],
        "dist": data["dist"],
        "n_alertas": data["n_alertas"],
        "real_send": data["real_send"],
        "send_flag": data["send_flag"],
        "email_enabled": data["email_on"],
        "telegram_enabled": data["tg_on"],
        "hist_note": data["hist_note"],
        "email_subject": data["email_subject"],
        "telegram_preview": data["tg_body"],
        "email_preview_markdown": data["email_body"],
        "legacy_level_change": {"subject": data["legacy_subject"], "body": data["legacy_body"]},
        "estadual": {
            "titulo": est.get("titulo"),
            "nivel": est.get("nivel"),
            "motivo": est.get("motivo"),
            "orientacoes": est.get("orientacoes"),
            "predicao": est.get("predicao"),
        },
        "cuiaba": {
            "titulo": cui.get("titulo"),
            "nivel": cui.get("nivel"),
            "motivo": cui.get("motivo"),
        },
    }
    side_path = out.with_suffix(".alertas_preview.json")
    side_path.write_text(json.dumps(side, ensure_ascii=False, indent=2), encoding="utf-8")
    # também txt para leitura rápida
    txt_path = OUT_DIR / "PREVIA_ALERTA_TELEGRAM_EMAIL.txt"
    txt_path.write_text(
        "=== PRÉVIA — NÃO ENVIADO (SEND_ALERT_ON_LEVEL_CHANGE=false) ===\n\n"
        f"ASSUNTO E-MAIL:\n{data['email_subject']}\n\n"
        f"CORPO TELEGRAM/E-MAIL (markdown):\n{data['md_est']}\n\n"
        f"--- CUIABÁ ---\n{data['md_cui']}\n",
        encoding="utf-8",
    )
    return out, data


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=str(DEFAULT_OUT))
    ap.add_argument("--also-downloads", action="store_true")
    args = ap.parse_args()
    path, data = build(Path(args.out))
    print(f"[OK] {path}")
    print(f"[OK] preview JSON: {path.with_suffix('.alertas_preview.json')}")
    print(f"[INFO] real_send={data['real_send']} · SEND_ALERT={data['send_flag']} · "
          f"email={data['email_on']} tg={data['tg_on']} · estadual={data['est'].get('nivel')}")
    if args.also_downloads:
        dest = Path.home() / "Downloads" / path.name
        shutil.copy2(path, dest)
        print(f"[OK] cópia: {dest}")
        shutil.copy2(OUT_DIR / "PREVIA_ALERTA_TELEGRAM_EMAIL.txt", Path.home() / "Downloads" / "PREVIA_ALERTA_TELEGRAM_EMAIL.txt")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
